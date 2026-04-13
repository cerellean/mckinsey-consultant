#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
McKinsey PPT Generator V4 - 示例代码
展示如何生成符合McKinsey标准的PPT
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

# McKinsey配色常量
PRIMARY_BLUE = RGBColor(0, 41, 96)
SECONDARY_BLUE = RGBColor(0, 101, 189)
LIGHT_BLUE = RGBColor(201, 240, 255)
IKEA_YELLOW = RGBColor(255, 219, 0)
GRAY = RGBColor(128, 128, 128)
LIGHT_GRAY = RGBColor(245, 245, 245)
WHITE = RGBColor(255, 255, 255)

# 布局常量
CONTENT_WIDTH = 8.4
CONTENT_LEFT = 0.8

class McKinseyPPTGenerator:
    """McKinsey风格PPT生成器"""
    
    def __init__(self):
        self.prs = Presentation()
        self.prs.slide_width = Inches(10)
        self.prs.slide_height = Inches(7.5)
    
    def create_content_box(self, slide, x, y, width, height, background=LIGHT_GRAY):
        """
        创建McKinsey风格内容框
        - 直角矩形（大文本框不用圆角）
        - 无边框
        - 背景色区分
        """
        box = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,  # ⭐ 直角矩形，不是ROUNDED_RECTANGLE
            Inches(x), Inches(y),
            Inches(width), Inches(height)
        )
        box.fill.solid()
        box.fill.fore_color.rgb = background
        box.line.fill.background()  # ⭐ 无边框
        return box
    
    def create_title(self, slide, text):
        """创建标题（Storyline驱动）"""
        title_box = slide.shapes.add_textbox(
            Inches(CONTENT_LEFT), 
            Inches(0.45), 
            Inches(CONTENT_WIDTH), 
            Inches(0.6)
        )
        title_box.text_frame.word_wrap = True
        p = title_box.text_frame.paragraphs[0]
        p.text = text
        p.font.size = Pt(20)
        p.font.bold = True
        p.font.color.rgb = PRIMARY_BLUE
        title_box.line.fill.background()
        return title_box
    
    def create_insight_box(self, slide, y, text):
        """
        创建洞察框
        - 直角矩形
        - 功能性边框（突出重要信息）
        """
        box = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(CONTENT_LEFT), 
            Inches(y),
            Inches(CONTENT_WIDTH), 
            Inches(0.6)
        )
        box.fill.solid()
        box.fill.fore_color.rgb = LIGHT_BLUE
        box.line.color.rgb = SECONDARY_BLUE
        box.line.width = Pt(1.5)  # 功能性边框
        
        text_box = slide.shapes.add_textbox(
            Inches(CONTENT_LEFT + 0.2), 
            Inches(y + 0.12), 
            Inches(CONTENT_WIDTH - 0.4), 
            Inches(0.36)
        )
        text_box.text_frame.word_wrap = True
        p = text_box.text_frame.paragraphs[0]
        p.text = text
        p.font.size = Pt(11)
        p.font.color.rgb = PRIMARY_BLUE
        text_box.line.fill.background()
        
        return box
    
    def create_table_with_header(self, slide, x, y, headers, data):
        """
        创建带深蓝表头的表格
        - 表头背景：深蓝
        - 表头文字：白色（强制）
        """
        col_widths = [2.0, 2.0, 2.0, 2.0]
        
        # 表头
        for i, header in enumerate(headers):
            header_box = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE,
                Inches(x + sum(col_widths[:i])), 
                Inches(y),
                Inches(col_widths[i]), 
                Inches(0.35)
            )
            header_box.fill.solid()
            header_box.fill.fore_color.rgb = PRIMARY_BLUE
            header_box.line.fill.background()
            
            text_box = slide.shapes.add_textbox(
                Inches(x + sum(col_widths[:i]) + 0.05), 
                Inches(y + 0.07),
                Inches(col_widths[i] - 0.1), 
                Inches(0.21)
            )
            text_box.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
            p = text_box.text_frame.paragraphs[0]
            p.text = header
            p.font.size = Pt(10)
            p.font.bold = True
            p.font.color.rgb = WHITE  # ⭐ 深蓝背景强制白色文字
            p.alignment = PP_ALIGN.CENTER
            text_box.line.fill.background()
        
        # 数据行
        for row_idx, row_data in enumerate(data):
            row_y = y + 0.4 + row_idx * 0.6
            for col_idx, cell in enumerate(row_data):
                cell_box = slide.shapes.add_shape(
                    MSO_SHAPE.RECTANGLE,
                    Inches(x + sum(col_widths[:col_idx])), 
                    Inches(row_y),
                    Inches(col_widths[col_idx]), 
                    Inches(0.55)
                )
                cell_box.fill.solid()
                cell_box.fill.fore_color.rgb = LIGHT_GRAY if col_idx == 0 else WHITE
                cell_box.line.color.rgb = GRAY
                cell_box.line.width = Pt(0.5)
                
                text_box = slide.shapes.add_textbox(
                    Inches(x + sum(col_widths[:col_idx]) + 0.05), 
                    Inches(row_y + 0.05),
                    Inches(col_widths[col_idx] - 0.1), 
                    Inches(0.45)
                )
                text_box.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
                p = text_box.text_frame.paragraphs[0]
                p.text = cell
                p.font.size = Pt(9)
                p.font.color.rgb = PRIMARY_BLUE
                p.alignment = PP_ALIGN.LEFT if col_idx == 0 else PP_ALIGN.CENTER
                text_box.line.fill.background()
    
    def add_footer(self, slide, center_text, page_num):
        """添加页脚"""
        left = slide.shapes.add_textbox(Inches(0.3), Inches(7.15), Inches(2.5), Inches(0.25))
        left.text_frame.text = "McKinsey & Company | Confidential"
        left.text_frame.paragraphs[0].font.size = Pt(8)
        left.text_frame.paragraphs[0].font.color.rgb = GRAY
        left.line.fill.background()
        
        center = slide.shapes.add_textbox(Inches(3.5), Inches(7.15), Inches(3), Inches(0.25))
        center.text_frame.text = center_text
        center.text_frame.paragraphs[0].font.size = Pt(8)
        center.text_frame.paragraphs[0].font.color.rgb = GRAY
        center.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        center.line.fill.background()
        
        right = slide.shapes.add_textbox(Inches(9.1), Inches(7.15), Inches(0.6), Inches(0.25))
        right.text_frame.text = str(page_num)
        right.text_frame.paragraphs[0].font.size = Pt(8)
        right.text_frame.paragraphs[0].font.color.rgb = GRAY
        right.text_frame.paragraphs[0].alignment = PP_ALIGN.RIGHT
        right.line.fill.background()
    
    def enforce_text_contrast(self, shape):
        """强制深蓝背景使用白色文字"""
        if not shape.has_text_frame:
            return
        
        try:
            if shape.fill.type == 1:
                bg = shape.fill.fore_color.rgb
                if bg == PRIMARY_BLUE or bg == SECONDARY_BLUE:
                    for para in shape.text_frame.paragraphs:
                        for run in para.runs:
                            run.font.color.rgb = WHITE
        except:
            pass
    
    def final_quality_check(self):
        """生成后质量检查"""
        print("\n执行最终质量检查...")
        
        for slide_idx, slide in enumerate(self.prs.slides):
            for shape in slide.shapes:
                # 检查颜色对比度
                self.enforce_text_contrast(shape)
                
                # 检查大文本框是否误用圆角
                if shape.has_text_frame:
                    width = shape.width / 914400
                    if width > 3 and hasattr(shape, 'shape_type'):
                        if 'ROUNDED' in str(shape.shape_type):
                            print(f"⚠️  第{slide_idx+1}页: {width:.1f}英寸宽的大文本框使用了圆角")
        
        print("✅ 质量检查完成")
    
    def save(self, filename):
        """保存PPT"""
        self.final_quality_check()
        self.prs.save(filename)
        print(f"\n✅ PPT已保存: {filename}")


# 使用示例
if __name__ == "__main__":
    gen = McKinseyPPTGenerator()
    
    # 添加封面
    slide = gen.prs.slides.add_slide(gen.prs.slide_layouts[6])
    gen.create_title(slide, "McKinsey风格PPT示例")
    
    # 添加内容页
    slide = gen.prs.slides.add_slide(gen.prs.slide_layouts[6])
    gen.create_title(slide, "三大趋势正在重塑行业格局，企业需要立即行动")
    
    # 内容框（直角矩形，无边框）
    gen.create_content_box(slide, CONTENT_LEFT, 1.3, CONTENT_WIDTH, 2.5)
    
    # 洞察框
    gen.create_insight_box(slide, 4.2, 
        "核心洞察：数字化转型是生存题而非选择题，窗口期仅剩3年")
    
    gen.add_footer(slide, "战略概览", 2)
    
    # 保存
    gen.save("mckinsey_example.pptx")
